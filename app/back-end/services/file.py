import csv
import mimetypes
import os
from io import BufferedReader
from typing import Optional

import docx2txt
import fitz
import pptx
from const import PAGE_SPLITTER
from fastapi import UploadFile
from loguru import logger
from pydantic_schemas.document import Document, DocumentMetadata

# import tabula


async def get_document_from_file(
    file: UploadFile, metadata: DocumentMetadata
) -> Document:
    extracted_text = await extract_text_from_form_file(file)

    doc = Document(text=extracted_text, metadata=metadata)

    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    try:
        with open(filepath, "rb") as file:
            extracted_text = extract_text_from_file(file, mimetype)
    except Exception as e:
        logger.error(e)
        raise e

    return extracted_text


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        pages = fitz.open(file)
        extracted_text = ""
        for page_num, page in enumerate(pages):
            # tables_on_page = tabula.read_pdf(
            #     file, pages=page_num + 1, multiple_tables=True
            # )
            # get_all text in the page
            page_text = page.get_text("text", sort=True)
            # convert table to html
            # if len(tables_on_page) > 0:
            #     tables_markdown = [table.to_html() for table in tables_on_page]
            #     extracted_text = page_text + "\n".join(tables_markdown)

            extracted_text = extracted_text + "\n" + page_text + PAGE_SPLITTER
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        raise ValueError("Unsupported file type: {}".format(mimetype))

    return extracted_text


# Extract text from a file based on its mimetype
async def extract_text_from_form_file(file: UploadFile):
    """Return the text content of a file."""
    # get the file body from the upload file object
    mimetype = file.content_type
    logger.info(f"mimetype: {mimetype}")
    logger.info(f"file.file: {file.file}")
    logger.info("file: ", file)

    file_stream = await file.read()

    temp_file_path = "/tmp/temp_file"

    # write the file to a temporary location
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)

    try:
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype)
    except Exception as e:
        logger.error(e)
        os.remove(temp_file_path)
        raise e

    # remove file from temp location
    os.remove(temp_file_path)

    return extracted_text
